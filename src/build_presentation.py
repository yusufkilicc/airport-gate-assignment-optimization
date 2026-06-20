from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(27, 62, 84)
TEAL = RGBColor(31, 94, 122)
RUST = RGBColor(204, 90, 60)
GOLD = RGBColor(215, 140, 31)
DARK = RGBColor(35, 35, 35)
LIGHT = RGBColor(244, 246, 248)
WHITE = RGBColor(255, 255, 255)


def add_banner(slide, title: str, subtitle: str | None = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9.2), Inches(0.35))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.52), Inches(0.50), Inches(8.4), Inches(0.22))
        sub_frame = sub_box.text_frame
        sp = sub_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(10.5)
        sp.font.color.rgb = WHITE


def add_bullets(
    slide,
    bullets: list[str],
    left: float = 0.7,
    top: float = 1.3,
    width: float = 5.8,
    height: float = 4.8,
    font_size: float = 20,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    first = True
    for bullet in bullets:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(7)
        first = False


def add_metric_card(slide, title: str, value: str, left: float, top: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(2.35),
        Inches(1.4),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.18), Inches(2.0), Inches(0.35))
    title_frame = title_box.text_frame
    tp = title_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(12)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    value_box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.58), Inches(2.0), Inches(0.5))
    value_frame = value_box.text_frame
    vp = value_frame.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(24)
    vp.font.bold = True
    vp.font.color.rgb = WHITE


def add_picture(slide, image_path: Path, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def add_caption(slide, text: str, left: float, top: float, width: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    frame = box.text_frame
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER


def build_presentation() -> Path:
    root = Path(__file__).resolve().parents[1]
    figures = root / "outputs" / "figures"
    report_dir = root / "reports"

    kpi = pd.read_csv(root / "outputs" / "kpi_summary.csv").set_index("solution")
    sim = pd.read_csv(root / "outputs" / "simulation_summary.csv").set_index("method")

    base = kpi.loc["Baseline FCFS"]
    priority = kpi.loc["Priority Dispatch"]
    opt = kpi.loc["Optimized MILP"]
    sim_opt = sim.loc["Optimized MILP"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def new_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = LIGHT
        return slide

    slide = new_slide()
    add_banner(slide, "Airport Turnaround Capacity and Gate Allocation Analysis", "INDE 4313 - Supply Chain Management")
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(6.2), Inches(1.5))
    tf = title_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "A Supply Chain Perspective"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = "Semester project on gate assignment, congestion, robustness, and service-oriented optimization"
    p2.font.size = Pt(17)
    p2.font.color.rgb = DARK
    add_picture(slide, figures / "capacity_profile.png", left=7.05, top=1.35, width=5.65)

    slide = new_slide()
    add_banner(slide, "Problem And Case Setup")
    add_bullets(
        slide,
        [
            "Airports behave like service supply chains with synchronized flows of aircraft, passengers, and time.",
            "The case models 36 flights across three traffic banks with 7 contact gates and 4 remote stands.",
            "Peak occupied-stand demand reaches 11, so remote usage is structurally unavoidable.",
            "The planning question is not only whether remote usage exists, but which flights should bear it.",
        ],
        width=6.3,
        font_size=18,
    )
    add_metric_card(slide, "Flights", "36", 7.35, 1.45, TEAL)
    add_metric_card(slide, "Passengers", "7,290", 9.95, 1.45, NAVY)
    add_metric_card(slide, "Contact Gates", "7", 7.35, 3.1, RUST)
    add_metric_card(slide, "Peak Demand", "11", 9.95, 3.1, GOLD)

    slide = new_slide()
    add_banner(slide, "Policies And Model")
    add_bullets(
        slide,
        [
            "Baseline FCFS assigns flights in arrival order and picks the first feasible gate.",
            "Priority Dispatch processes high-passenger flights first to protect premium gates for higher-impact flights.",
            "Optimized MILP minimizes walking, taxi penalty, remote usage penalty, zone mismatch, and gate-scarcity misuse.",
            "Compatibility and overlap constraints ensure every flight gets exactly one feasible stand.",
        ],
        left=0.8,
        top=1.45,
        width=7.0,
        height=4.6,
        font_size=18,
    )
    add_metric_card(slide, "Model Type", "Binary MILP", 8.8, 1.6, NAVY)
    add_metric_card(slide, "Solver", "CBC", 8.8, 3.25, TEAL)
    add_metric_card(slide, "Methods", "3", 8.8, 4.9, GOLD)

    slide = new_slide()
    add_banner(slide, "Capacity Pressure")
    add_picture(slide, figures / "capacity_profile.png", left=0.55, top=1.2, width=8.15)
    add_bullets(
        slide,
        [
            "Repeated banked demand pushes occupied stands above contact-gate capacity.",
            "This makes the project a resource-allocation problem under unavoidable congestion.",
            "The managerial objective becomes better prioritization of scarce premium positions.",
        ],
        left=9.0,
        top=1.6,
        width=3.7,
        height=4.0,
        font_size=17,
    )

    slide = new_slide()
    add_banner(slide, "KPI Comparison Across Three Methods")
    add_picture(slide, figures / "kpi_comparison.png", left=0.45, top=1.15, width=8.2)
    add_bullets(
        slide,
        [
            f"Priority Dispatch reduces total cost from {base['total_assignment_cost']:.0f} to {priority['total_assignment_cost']:.0f}.",
            f"MILP lowers total cost further to {opt['total_assignment_cost']:.0f} and cuts zone mismatches from {int(base['zone_mismatch_flights'])} to {int(opt['zone_mismatch_flights'])}.",
            f"Remote passengers fall from 3014 in FCFS to 2149 in MILP even though remote-flight count stays at 11.",
            "Priority Dispatch is a strong low-complexity heuristic, but MILP is best overall.",
        ],
        left=8.9,
        top=1.3,
        width=3.8,
        height=5.2,
        font_size=16,
    )

    slide = new_slide()
    add_banner(slide, "Gate Schedule Comparison")
    add_picture(slide, figures / "baseline_gantt.png", left=0.25, top=1.25, width=4.2)
    add_picture(slide, figures / "priority_gantt.png", left=4.55, top=1.25, width=4.2)
    add_picture(slide, figures / "optimized_gantt.png", left=8.85, top=1.25, width=4.2)
    add_caption(slide, "Baseline FCFS", 1.2, 6.75, 2.4)
    add_caption(slide, "Priority Dispatch", 5.45, 6.75, 2.7)
    add_caption(slide, "Optimized MILP", 9.8, 6.75, 2.4)

    slide = new_slide()
    add_banner(slide, "Sensitivity Analysis")
    add_picture(slide, figures / "sensitivity_analysis.png", left=0.45, top=1.2, width=8.35)
    add_bullets(
        slide,
        [
            "Changing the remote-stand penalty alters total cost but not the number of remote flights in the tested range.",
            "That means remote usage is structural in this case, not just a penalty-calibration artifact.",
            "Zone and scarcity penalties change allocation behavior more visibly than the remote penalty.",
        ],
        left=8.95,
        top=1.45,
        width=3.55,
        height=4.9,
        font_size=16,
    )

    slide = new_slide()
    add_banner(slide, "Robustness Simulation")
    add_picture(slide, figures / "simulation_results.png", left=0.5, top=1.2, width=8.2)
    add_bullets(
        slide,
        [
            f"MILP has the lowest mean simulated cost at {sim_opt['total_cost_mean']:.0f}.",
            f"It also achieves the best feasibility rate at {sim_opt['feasibility_rate']:.0%}.",
            "Priority Dispatch stays consistently between MILP and FCFS under random delays.",
            "Optimization therefore improves both planned and disrupted-day performance.",
        ],
        left=8.95,
        top=1.35,
        width=3.6,
        height=5.0,
        font_size=16,
    )

    slide = new_slide()
    add_banner(slide, "Scenario Analysis")
    add_picture(slide, figures / "scenario_comparison.png", left=0.55, top=1.35, width=7.8)
    add_bullets(
        slide,
        [
            "Normal and disruption scenarios remain feasible and still favor MILP.",
            "The heavy scenario becomes infeasible under the current stand inventory.",
            "This turns scenario analysis into a capacity-warning tool, not only a KPI comparison.",
        ],
        left=8.75,
        top=1.55,
        width=3.8,
        height=4.6,
        font_size=17,
    )

    slide = new_slide()
    add_banner(slide, "Dashboard And Deliverables")
    add_bullets(
        slide,
        [
            "The project includes a Dash dashboard with four tabs: solution comparison, sensitivity, simulation, and scenario analysis.",
            "Static deliverables include the report, processed data, KPI tables, scenario outputs, and figure set.",
            "All analysis artifacts are reproducible from a single pipeline: python src/run_project.py",
        ],
        left=0.85,
        top=1.55,
        width=12.0,
        height=4.8,
        font_size=18,
    )

    slide = new_slide()
    add_banner(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "Gate assignment is a supply chain planning problem, not only a dispatch rule.",
            "Priority Dispatch provides a strong heuristic improvement over FCFS.",
            "MILP delivers the best overall service and robustness outcomes.",
            "Heavy-demand infeasibility shows that capacity limits can dominate policy quality.",
        ],
        left=0.9,
        top=1.55,
        width=7.0,
        height=4.2,
        font_size=18,
    )
    add_metric_card(slide, "MILP Cost Gain", "14.7%", 8.35, 1.55, NAVY)
    add_metric_card(slide, "Walk Reduction", "17.95%", 8.35, 3.2, TEAL)
    add_metric_card(slide, "Heavy Scenario", "Infeasible", 8.35, 4.85, RUST)

    output_path = report_dir / "INDE4313_presentation.pptx"
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build_presentation()
    print(f"Presentation created: {path}")
